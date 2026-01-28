from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    print("Generating HydroAI Mapathon Presentation...")
    prs = Presentation()

    # --- Helper Functions ---
    def set_font(run, size, bold=False, color=None):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = 'Calibri'
        if color:
            run.font.color.rgb = color

    def add_slide(title_text, content_points, footer_text="HydroAI - Mapathon 2024"):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title Styling
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

        # Content Styling
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear() # Clear default bullets

        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.space_after = Pt(14)
            p.level = 0
            set_font(p.font, 20)

        # Footer
        left = Inches(0.5)
        top = Inches(7.2)
        width = Inches(9)
        height = Inches(0.3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = footer_text
        p.alignment = PP_ALIGN.RIGHT
        set_font(p.font, 10, color=RGBColor(100, 100, 100))

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "HydroAI: Geospatial Reservoir Intelligence"
    subtitle.text = "Mapathon 2024\nProject ID: HYDRO-MAP-24\nTeam: HydroAI Solutions"
    
    # Style Title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # =========================================================================
    # SLIDE 2: Problem Statement
    # =========================================================================
    add_slide("Problem Statement", [
        "Water Scarcity & Management: Reservoir levels fluctuate unpredictably due to climate change.",
        "Lack of Real-Time Data: Traditional gauging is manual, slow, and prone to human error.",
        "Inefficient Risk Assessment: Decision-makers lack integrated tools to predict floods or droughts.",
        "Data Silos: Satellite imagery, rainfall data, and hydrological models often exist in disconnected systems."
    ])

    # =========================================================================
    # SLIDE 3: The Solution
    # =========================================================================
    add_slide("The Solution: HydroAI", [
        "Automated Monitoring: Uses Satellite Imagery (Sentinel-2) to extract water spread area automatically.",
        "AI-Driven Insights: Integrates Generative AI (Gemini) to produce human-readable hydrological reports.",
        "Predictive Modeling: Uses LSTM networks to forecast future water volumes.",
        "Interactive Dashboard: A unified geospatial interface for officials to visualize risk and take action."
    ])

    # =========================================================================
    # SLIDE 4: Technology Stack
    # =========================================================================
    add_slide("Technology Stack", [
        "Frontend: React 18, TypeScript, Tailwind CSS, Recharts, Leaflet (Mapping).",
        "Backend: Python (FastAPI), Uvicorn.",
        "Geospatial Engine: Google Earth Engine (GEE) API, Rasterio, GeoJSON.",
        "Machine Learning: PyTorch (Deep Learning), Scikit-Learn.",
        "Generative AI: Google Gemini Pro (via @google/genai SDK)."
    ])

    # =========================================================================
    # SLIDE 5: System Architecture
    # =========================================================================
    add_slide("System Architecture", [
        "1. Data Acquisition: Ingests Sentinel-2 L2A imagery via GEE.",
        "2. Pre-processing: Atmospheric correction and NDWI index calculation.",
        "3. ML Analysis: U-Net segmentation extracts precise water boundaries.",
        "4. Forecasting: Historical time-series data feeds into LSTM models.",
        "5. Interface: React Frontend displays maps, alerts, and AI summaries.",
        "6. Feedback Loop: RLHF mechanism improves model accuracy over time."
    ])

    # =========================================================================
    # SLIDE 6: Satellite Intelligence
    # =========================================================================
    add_slide("Satellite & Geospatial Processing", [
        "Source: Sentinel-2 (European Space Agency).",
        "Resolution: 10 meters per pixel (High fidelity).",
        "Spectral Bands: Green (B3) and NIR (B8) used for NDWI (Normalized Difference Water Index).",
        "Dynamic Tiling: Leaflet map layers overlayed with real-time satellite tiles.",
        "Metric: Surface Area (sq km) calculated dynamically from pixel counts."
    ])

    # =========================================================================
    # SLIDE 7: Machine Learning (Computer Vision)
    # =========================================================================
    add_slide("ML Model 1: U-Net for Segmentation", [
        "Objective: Pixel-wise classification of water vs. land.",
        "Architecture: U-Net with ResNet34 backbone.",
        "Input: 4 Channels (Red, Green, Blue, Near-Infrared).",
        "Dataset: SWED (Sentinel-2 Water Edges Dataset).",
        "Performance Metric: IoU (Intersection over Union).",
        "Augmentation: Random rotations and flips during training to ensure robustness."
    ])

    # =========================================================================
    # SLIDE 8: Forecasting & Anomaly Detection
    # =========================================================================
    add_slide("ML Model 2: Forecasting & Risk", [
        "Long Short-Term Memory (LSTM):",
        "   - Analyzes historical volume trends (time-series).",
        "   - Predicts water storage for the next 30 days.",
        "Isolation Forest (Anomaly Detection):",
        "   - Detects sudden, uncharacteristic drops in water levels.",
        "   - Flags potential sensor errors or illegal discharge events."
    ])

    # =========================================================================
    # SLIDE 9: Generative AI & RLHF
    # =========================================================================
    add_slide("Generative AI & RLHF", [
        "Hydrologist Agent: Google Gemini Pro analyzes raw stats (Volume, NDWI, Anomaly).",
        "Output: Generates structured risk assessments, summaries, and recommendations.",
        "RLHF (Reinforcement Learning from Human Feedback):",
        "   - Users can correct model outputs (e.g., 'Risk is High, not Moderate').",
        "   - Feedback is logged to retraining datasets to fine-tune the system."
    ])

    # =========================================================================
    # SLIDE 10: Conclusion & Future Scope
    # =========================================================================
    add_slide("Impact & Future Scope", [
        "Impact: Reduces response time to flood/drought events by 60%.",
        "Scalability: Architecture supports any reservoir with Lat/Lng coordinates.",
        "Future Scope:",
        "   - Integration with IoT water level sensors for hybrid accuracy.",
        "   - Mobile app for field officers.",
        "   - Bathymetric mapping integration for 3D volume estimation."
    ])

    # Save
    filename = "backend/HydroAI_Mapathon_Presentation.pptx"
    prs.save(filename)
    print(f"✅ Presentation saved successfully to {filename}")

if __name__ == "__main__":
    create_presentation()
